from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# 投影片內容（精簡成10分鐘版本）
slides = [
    ("FastAPI 介紹", "現代 Python Web Framework\n高效能 × 易用性 × 自動文件"),
    
    ("什麼是 FastAPI", 
     "• Python API 框架\n• 基於 Starlette + Pydantic\n• 高效能設計"),
    
    ("為什麼用 FastAPI",
     "• 🚀 高效能\n• 📄 自動文件\n• ⚡ 開發快速\n• 🔒 資料驗證"),
    
    ("安裝與啟動",
     "pip install fastapi uvicorn\nuvicorn main:app --reload"),
    
    ("Hello World",
     "from fastapi import FastAPI\napp = FastAPI()\n@app.get('/')\ndef read_root():\n    return {'Hello': 'World'}"),
    
    ("路由與參數",
     "• 支援 path / query\n• 自動型別檢查\n• 例：/items/{id}"),
    
    ("資料模型",
     "使用 Pydantic\n• 自動驗證\n• 清楚定義 API 結構"),
    
    ("自動文件",
     "• /docs (Swagger UI)\n• /redoc\n• 可直接測試 API"),
    
    ("非同步支援",
     "async / await\n提升 I/O 效能\n適合高併發"),
    
    ("FastAPI vs Flask",
     "FastAPI：快 + 自動文件"),
    
    ("應用場景",
     "• REST API\n• AI / ML 後端\n• 微服務\n• 即時應用"),
    
    ("總結",
     "• 高效能\n• 開發快\n• 現代 API 首選"),
    
    ("Q&A", "Thank You 🙌")
]

for title, content in slides:
    add_slide(title, content)

prs.save("FastAPI_10min_Presentation.pptx")